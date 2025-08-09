import os, httpx, requests, random, difflib
from pptx import Presentation
from pptx.util import Inches

TEMPLATE_DIR = "templates"
PERPLEXITY_API_KEY = os.getenv("PERPLEXITY_API_KEY")
IMAGE_GEN_MODEL = "sonar-pro"

IMPORTANT_SLIDE_TYPES = {"Problem", "Solution", "Product"}  # Slide types eligible for image generation


def generate_image_via_perplexity(prompt: str) -> str:
    """
    Calls Perplexity API to generate an image based on the prompt.
    Returns the image URL (if available).
    """
    try:
        resp = httpx.post(
            "https://api.perplexity.ai/chat/completions",
            headers={"Authorization": f"Bearer {PERPLEXITY_API_KEY}"},
            json={
                "model": IMAGE_GEN_MODEL,
                "stream": False,
                "messages": [
                    {"role": "user", "content": f"Generate a realistic, clean, business-style illustration of: {prompt}"}
                ]
            },
            timeout=30
        )
        resp.raise_for_status()
        data = resp.json()
        # You'll need to adapt this once Perplexity confirms the actual image key structure
        return data["choices"][0]["message"].get("image_url")  # Placeholder path
    except Exception as e:
        print(f"Perplexity API Error: {e}")
        return None


def download_image(url: str, dest_folder="images", filename=None) -> str:
    """
    Download the image from URL and return the local path.
    """
    if not os.path.exists(dest_folder):
        os.makedirs(dest_folder)
    if not filename:
        filename = f"img_{int(time.time())}.jpg"
    path = os.path.join(dest_folder, filename)

    try:
        response = requests.get(url)
        response.raise_for_status()
        with open(path, "wb") as f:
            f.write(response.content)
        return path
    except Exception as e:
        print(f"Image download failed: {e}")
        return None


def add_image_to_slide(slide, image_path, left=Inches(6), top=Inches(1.5), width=Inches(3)):
    """
    Add image to the given slide.
    """
    try:
        slide.shapes.add_picture(image_path, left, top, width=width)
    except Exception as e:
        print(f"Failed to add image to slide: {e}")


def save_as_powerpoint(pitch_deck: dict,
                       output_filename: str = "pitch_deck.pptx",
                       industry: str = "Generic") -> None:
    # 1. Gather all templates
    all_templates = [
        os.path.join(TEMPLATE_DIR, f)
        for f in os.listdir(TEMPLATE_DIR)
        if f.lower().endswith(".pptx")
    ]
    if not all_templates:
        raise RuntimeError(f"No .pptx templates found in {TEMPLATE_DIR!r}")

    # 2. Match industry template
    industry_lc = industry.lower()
    direct_matches = [path for path in all_templates if industry_lc in os.path.basename(path).lower()]

    if not direct_matches:
        basenames = {os.path.splitext(os.path.basename(path))[0].lower(): path for path in all_templates}
        closest = difflib.get_close_matches(industry_lc, basenames.keys(), n=1, cutoff=0.6)
        if closest:
            direct_matches = [basenames[closest[0]]]

    # 3. Final fallback
    chosen_template = random.choice(direct_matches if direct_matches else all_templates)
    prs = Presentation(chosen_template)

    # 4. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = pitch_deck.get("title", "Pitch Deck")

    startup_name = pitch_deck.get('startup_info', {}).get('name', 'Unknown Startup')
    if len(slide.placeholders) > 1 and slide.placeholders[1]:
        slide.placeholders[1].text = f"For: {startup_name}\nIndustry: {industry}"

    # 5. Content Slides
    for i, slide_data in enumerate(pitch_deck.get("slides", [])):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data.get("title", "Untitled Slide")

        # Content bullets
        if len(slide.placeholders) > 1 and slide.placeholders[1] and hasattr(slide.placeholders[1], 'text_frame'):
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for bullet in slide_data.get("content", []):
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

        # # Image generation for key slides
        # if slide_data.get("type") in IMPORTANT_SLIDE_TYPES:
        #     prompt = f"{slide_data.get('title')} for a {industry} startup pitch deck"
        #     print(f"Generating image for slide '{slide_data['title']}' using Perplexity: {prompt}")
        #     image_url = generate_image_via_perplexity(prompt)
        #     if image_url:
        #         image_path = download_image(image_url, filename=f"slide_{i}.jpg")
        #         if image_path:
        #             add_image_to_slide(slide, image_path)
        # Add generated visuals directly into slide
        if slide_data.get("generated_visuals"):
            for img_path in slide_data["generated_visuals"]:
                add_image_to_slide(slide, img_path)

        # Notes with visuals (optional)
        visuals = slide_data.get("visual_elements", [])
        if visuals:
            try:
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_tf = slide.notes_slide.notes_text_frame
                    notes_tf.text = "Suggested Visuals:"
                    for vis in visuals:
                        p = notes_tf.add_paragraph()
                        p.text = f"• {vis}"
            except (AttributeError, TypeError):
                print(f"Warning: Could not add visuals to notes for slide '{slide_data.get('title', 'Untitled')}'")

    prs.save(output_filename)
    print(f"\n✅ Saved {output_filename} using template '{os.path.basename(chosen_template)}'")
